"""
Multi-format document parser.
Supports: PDF, DOCX, XLSX, PPTX, TXT.
"""

import logging
from pathlib import Path

logger = logging.getLogger(__name__)


def _parse_pdf(path: Path) -> str:
    try:
        import pdfplumber
        text_parts: list[str] = []
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                t = page.extract_text()
                if t:
                    text_parts.append(t)
        return "\n\n".join(text_parts)
    except Exception:
        # Fallback to PyMuPDF
        try:
            import fitz  # PyMuPDF
            doc = fitz.open(str(path))
            return "\n\n".join(page.get_text() for page in doc)
        except Exception as exc:
            logger.warning("PDF parse failed for %s: %s", path.name, exc)
            return ""


def _parse_docx(path: Path) -> str:
    try:
        from docx import Document
        doc = Document(str(path))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        tables: list[str] = []
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                if row_text.strip():
                    tables.append(row_text)
        return "\n".join(paragraphs) + ("\n\nTables:\n" + "\n".join(tables) if tables else "")
    except Exception as exc:
        logger.warning("DOCX parse failed for %s: %s", path.name, exc)
        return ""


def _parse_xlsx(path: Path) -> str:
    try:
        import pandas as pd
        xl = pd.ExcelFile(str(path))
        parts: list[str] = []
        for sheet_name in xl.sheet_names:
            df = xl.parse(sheet_name).fillna("")
            text = f"[Sheet: {sheet_name}]\n{df.to_string(index=False)}"
            parts.append(text)
        return "\n\n".join(parts)
    except Exception as exc:
        logger.warning("XLSX parse failed for %s: %s", path.name, exc)
        return ""


def _parse_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        slides: list[str] = []
        for i, slide in enumerate(prs.slides, 1):
            texts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                slides.append(f"[Slide {i}]\n" + "\n".join(texts))
        return "\n\n".join(slides)
    except Exception as exc:
        logger.warning("PPTX parse failed for %s: %s", path.name, exc)
        return ""


def _parse_txt(path: Path) -> str:
    for enc in ("utf-8", "utf-8-sig", "shift_jis", "cp932", "latin-1"):
        try:
            return path.read_text(encoding=enc)
        except (UnicodeDecodeError, LookupError):
            continue
    return ""


_PARSERS = {
    ".pdf": _parse_pdf,
    ".docx": _parse_docx,
    ".xlsx": _parse_xlsx,
    ".pptx": _parse_pptx,
    ".txt": _parse_txt,
}


def parse_file(path: Path) -> str:
    """Parse a document and return its plain-text content."""
    ext = path.suffix.lower()
    parser = _PARSERS.get(ext)
    if not parser:
        logger.debug("No parser for extension: %s", ext)
        return ""
    return parser(path)
